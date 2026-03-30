import os
os.environ["SCARF_NO_ANALYTICS"] = "true"
import shutil
import base64
import io
import traceback

from fastapi import UploadFile, HTTPException, BackgroundTasks
from langchain_core.messages import HumanMessage
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma

from core.config import settings
from core.chroma import chroma_client
from core.llm import embedding_function, llm_vision

from youtube_transcript_api import YouTubeTranscriptApi
import requests
import re


def _describe_image(b64_data: str, mime: str = "jpeg") -> str:
    """Call Vision LLM to describe a base64-encoded image. Returns description or empty string."""
    try:
        msg = HumanMessage(
            content=[
                {
                    "type": "text",
                    "text": (
                        "Describe this educational image, chart, or slide in extreme detail. "
                        "Focus heavily on extracting the core educational concepts, data points, "
                        "and semantic relationships shown. Ignore purely decorative visual elements."
                    ),
                },
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/{mime};base64,{b64_data}"},
                },
            ]
        )
        response = llm_vision.invoke([msg])
        return response.content.strip()
    except Exception as e:
        print(f"WARNING: Vision LLM failed: {e}")
        return ""


def _parse_pdf(path: str) -> list[Document]:
    """Extract text (and optionally image descriptions) from a PDF using pypdfium2."""
    import pypdfium2 as pdfium

    docs: list[Document] = []
    pdf = pdfium.PdfDocument(path)
    for page_num in range(len(pdf)):
        page = pdf[page_num]

        page_text = page.get_textpage().get_text_range().strip()
        
        page_text = re.sub(r'(?<![.:!?\n])\n(?!\n)', ' ', page_text)
        page_text = re.sub(r' {2,}', ' ', page_text)

        try:
            if len(page_text) < 200:
                bitmap = page.render(scale=150 / 72)
                pil_img = bitmap.to_pil()
                buf = io.BytesIO()
                pil_img.save(buf, format="JPEG", quality=85)
                b64 = base64.b64encode(buf.getvalue()).decode()
                
                vision_description = _describe_image(b64, "jpeg")
                if vision_description:
                    page_text += f"\n\n[Visual Element Description: {vision_description}]"
        except Exception as e:
            print(f"WARNING: PDF image render failed (page {page_num + 1}): {e}")

        page_text = page_text.strip()
        if page_text:
            docs.append(Document(page_content=page_text, metadata={"page": page_num + 1}))

    pdf.close()
    return docs


def _parse_pptx(path: str) -> list[Document]:
    """Extract text and image descriptions from a PPTX file."""
    from pptx import Presentation

    docs: list[Document] = []
    prs = Presentation(path)

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)

            if shape.shape_type == 13:
                try:
                    img_bytes = shape.image.blob
                    b64 = base64.b64encode(img_bytes).decode()
                    ext = shape.image.ext.lower()
                    mime = "jpeg" if ext in ("jpg", "jpeg") else ext
                    desc = _describe_image(b64, mime)
                    if desc:
                        slide_texts.append(f"[Image Description: {desc}]")
                except Exception as e:
                    print(f"WARNING: PPTX image extraction failed (slide {slide_num}): {e}")

        if slide_texts:
            docs.append(Document(
                page_content="\n".join(slide_texts),
                metadata={"slide": slide_num},
            ))

    return docs


def _parse_docx(path: str) -> list[Document]:
    """Extract text and inline image descriptions from a DOCX file."""
    from docx import Document as DocxDocument
    from docx.oxml.ns import qn

    docx = DocxDocument(path)
    docs: list[Document] = []

    all_text: list[str] = []
    for para in docx.paragraphs:
        text = para.text.strip()
        if text:
            all_text.append(text)

        for run in para.runs:
            for drawing in run._element.findall(f".//{qn('a:blip')}", namespaces=run._element.nsmap):
                try:
                    rId = drawing.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
                    if rId:
                        image_part = docx.part.related_parts.get(rId)
                        if image_part:
                            img_bytes = image_part.blob
                            b64 = base64.b64encode(img_bytes).decode()
                            desc = _describe_image(b64, "jpeg")
                            if desc:
                                all_text.append(f"[Image Description: {desc}]")
                except Exception as e:
                    print(f"WARNING: DOCX image failed: {e}")

    for table in docx.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                all_text.append(row_text)

    if all_text:
        docs.append(Document(page_content="\n".join(all_text)))

    return docs


def _parse_text(path: str) -> list[Document]:
    """Parse plain text / markdown files."""
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()
    return [Document(page_content=content)] if content.strip() else []


def extract_video_id(url: str) -> str:
    regex = r"(?:youtube\.com\/(?:[^\/]+\/.+\/|(?:v|e(?:mbed)?)\/|.*[?&]v=)|youtu\.be\/)([^\"&?\/\s]{11})"
    match = re.search(regex, url)
    if match:
        return match.group(1)
    raise ValueError("Could not extract a valid 11-character YouTube Video ID from the provided URL.")

def _parse_youtube(url: str) -> tuple[str, list[Document]]:
    """Parse a YouTube URL to extract its title and time-stamped transcript."""
    video_id = extract_video_id(url)

    title = f"YouTube Video ({video_id})"
    try:
        oembed_url = f"https://www.youtube.com/oembed?url={url}&format=json"
        response = requests.get(oembed_url, timeout=5)
        if response.status_code == 200:
            title = response.json().get("title", title)
    except Exception as e:
        print(f"WARNING: oEmbed fetch failed: {e}")

    try:
        ytt_api = YouTubeTranscriptApi()
        transcript = ytt_api.fetch(video_id)
    except Exception as e:
        raise ValueError(f"Could not retrieve transcript. Ensure the video has closed captions enabled. (Details: {e})")

    text_parts = []
    for snippet in transcript:
        text = getattr(snippet, 'text', None) or snippet.get('text', '')
        start_seconds = getattr(snippet, 'start', None) or snippet.get('start', 0.0)
        
        m, s = divmod(int(start_seconds), 60)
        h, m = divmod(m, 60)
        timestamp = f"**[{h:02d}:{m:02d}:{s:02d}]**" if h > 0 else f"**[{m:02d}:{s:02d}]**"
        text_parts.append(f"{timestamp} {text}")

    full_text = "\n".join(text_parts)
    return title, [Document(page_content=full_text, metadata={"source": url, "title": title})]



async def ingest_document(
    file: UploadFile,
    doc_id: str,
    user_id: str,
    background_tasks: BackgroundTasks,
) -> dict:
    """
    Save the uploaded file, chunk & embed it into ChromaDB, then queue the
    knowledge-graph + podcast pipeline as a background task.
    """
    os.makedirs(settings.TEMP_DIR, exist_ok=True)
    temp_path = os.path.join(settings.TEMP_DIR, file.filename)

    try:
        with open(temp_path, "wb") as buf:
            shutil.copyfileobj(file.file, buf)

        fname = file.filename.lower()
        print(f"DEBUG: Parsing '{fname}'...")

        if fname.endswith(".pdf"):
            documents = _parse_pdf(temp_path)
        elif fname.endswith((".pptx", ".ppt")):
            documents = _parse_pptx(temp_path)
        elif fname.endswith((".docx", ".doc")):
            documents = _parse_docx(temp_path)
        elif fname.endswith((".txt", ".md", ".markdown")):
            documents = _parse_text(temp_path)
        else:
            raise HTTPException(
                status_code=415,
                detail=f"Unsupported file type: {fname}. Supported: PDF, DOCX, PPTX, TXT, MD",
            )

        total_chars = sum(len(d.page_content) for d in documents)
        print(f"DEBUG: Extracted {len(documents)} blocks, {total_chars} chars total.")

        if total_chars == 0:
            raise HTTPException(
                status_code=400,
                detail="Document appears to be empty or unrecognised format.",
            )

        if fname.endswith(".pdf"):
            splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)
        elif fname.endswith((".pptx", ".ppt")):
            splitter = RecursiveCharacterTextSplitter(chunk_size=2000, chunk_overlap=400)
        else:
            splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
            
        chunks = splitter.split_documents(documents)

        for chunk in chunks:
            chunk.metadata["doc_id"] = doc_id
            chunk.metadata["source"] = file.filename
            
            header = f"Source Document: {file.filename}"
            if "page" in chunk.metadata:
                header += f", Page: {chunk.metadata['page']}"
            elif "slide" in chunk.metadata:
                header += f", Slide: {chunk.metadata['slide']}"
            
            chunk.page_content = f"[{header}]\n---\n{chunk.page_content}"

        print(f"DEBUG: Embedding {len(chunks)} chunks…")
        Chroma.from_documents(
            documents=chunks,
            embedding=embedding_function,
            persist_directory=settings.CHROMA_PATH,
            collection_name="mindnexus_docs",
        )
        print("DEBUG: Embeddings stored.")

        from services.background_service import process_graph_and_podcast
        background_tasks.add_task(process_graph_and_podcast, doc_id, user_id, chunks)
        print(f"DEBUG: Background pipeline queued for doc_id={doc_id}")

        return {"status": "success", "chunks": len(chunks)}

    except HTTPException:
        raise
    except Exception as exc:
        print(f"ERROR [ingest]: {exc}")
        raise HTTPException(status_code=500, detail=str(exc))
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)


async def ingest_youtube(
    url: str,
    doc_id: str,
    user_id: str,
    background_tasks: BackgroundTasks,
) -> dict:
    """
    Fetch a YouTube transcript, chunk & embed it into ChromaDB,
    then queue the KG + podcast pipeline.
    """
    try:
        title, documents = _parse_youtube(url)

        total_chars = sum(len(d.page_content) for d in documents)
        if total_chars == 0:
            raise HTTPException(
                status_code=400,
                detail="YouTube transcript appears to be empty.",
            )

        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        chunks = splitter.split_documents(documents)

        for chunk in chunks:
            chunk.metadata["doc_id"] = doc_id
            chunk.metadata["source"] = url
            chunk.metadata["title"] = title
            
            header = f"YouTube Video Title: {title}"
            chunk.page_content = f"[{header}]\n---\n{chunk.page_content}"

        print(f"DEBUG: Embedding {len(chunks)} YouTube chunks…")
        Chroma.from_documents(
            documents=chunks,
            embedding=embedding_function,
            persist_directory=settings.CHROMA_PATH,
            collection_name="mindnexus_docs",
        )
        print("DEBUG: YouTube Embeddings stored.")

        print(f"DEBUG: Skipping Neo4j Background pipeline for YouTube doc_id={doc_id}")

        return {"status": "success", "chunks": len(chunks), "title": title}

    except ValueError as val_exc:
        raise HTTPException(status_code=400, detail=str(val_exc))
    except Exception as exc:
        print(f"ERROR [ingest_youtube]: {exc}")
        traceback.print_exc()
        raise HTTPException(status_code=400, detail=f"Failed to ingest YouTube video: {str(exc)}")

async def reindex_document(
    doc_id: str,
    user_id: str,
    background_tasks: BackgroundTasks,
) -> dict:
    """
    Re-run the knowledge graph and podcast pipeline for an existing document
    by retrieving its text chunks directly from ChromaDB.
    """
    try:
        from core.chroma import collection
        from langchain_core.documents import Document
        from services.background_service import process_graph_and_podcast
        
        results = collection.get(where={"doc_id": doc_id}, include=["documents", "metadatas"])
        if not results or not results.get("documents"):
            raise HTTPException(status_code=404, detail="No embedded text found for this document")

        chunks = []
        for doc_text, meta in zip(results["documents"], results["metadatas"]):
            chunks.append(Document(page_content=doc_text, metadata=meta))

        from core.neo4j import graph
        if graph:
            graph.query("MATCH (n {doc_id: $doc_id}) DETACH DELETE n", {"doc_id": doc_id})
            print(f"DEBUG [reindex]: Wiped old graph nodes for {doc_id}")
            
        import os
        from core.config import settings
        for ext in ("mp3", "json"):
            path = os.path.join(settings.AUDIO_DIR, f"podcast_{doc_id}.{ext}")
            if os.path.exists(path):
                os.remove(path)
                print(f"DEBUG [reindex]: Wiped old audio {path}")
                
        print(f"DEBUG [reindex]: Queued background pipeline for {doc_id}")
        background_tasks.add_task(process_graph_and_podcast, doc_id, user_id, chunks)
        
        return {"status": "success", "message": "Re-indexing started"}
    except HTTPException:
        raise
    except Exception as exc:
        print(f"ERROR [reindex]: {exc}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(exc))
